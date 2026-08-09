#!/usr/bin/env python3
"""Build a PowerPoint deck from naturally sorted slide images."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


DEFAULT_EXTENSIONS = (".png", ".jpg", ".jpeg")


def natural_key(path: Path) -> list[object]:
    parts = re.split(r"(\d+)", path.name.lower())
    return [int(part) if part.isdigit() else part for part in parts]


def collect_images(input_dir: Path, extensions: tuple[str, ...]) -> list[Path]:
    normalized = {ext.lower() if ext.startswith(".") else f".{ext.lower()}" for ext in extensions}
    return sorted(
        (path for path in input_dir.iterdir() if path.is_file() and path.suffix.lower() in normalized),
        key=natural_key,
    )


def add_background(slide, width: int, height: int, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()


def add_fitted_picture(slide, image_path: Path, slide_width: int, slide_height: int, fit: str) -> None:
    if fit == "stretch":
        slide.shapes.add_picture(str(image_path), 0, 0, width=slide_width, height=slide_height)
        return

    with Image.open(image_path) as image:
        image_width, image_height = image.size

    if image_width <= 0 or image_height <= 0:
        raise ValueError(f"Invalid image dimensions: {image_path}")

    image_ratio = image_width / image_height
    slide_ratio = slide_width / slide_height

    if fit == "contain":
        if image_ratio >= slide_ratio:
            width = slide_width
            height = int(slide_width / image_ratio)
            left = 0
            top = int((slide_height - height) / 2)
        else:
            height = slide_height
            width = int(slide_height * image_ratio)
            top = 0
            left = int((slide_width - width) / 2)
    else:
        if image_ratio >= slide_ratio:
            height = slide_height
            width = int(slide_height * image_ratio)
            top = 0
            left = int((slide_width - width) / 2)
        else:
            width = slide_width
            height = int(slide_width / image_ratio)
            left = 0
            top = int((slide_height - height) / 2)

    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def build_deck(args: argparse.Namespace) -> dict[str, object]:
    input_dir = args.input_dir.resolve()
    output = args.output.resolve()

    if not input_dir.is_dir():
        raise FileNotFoundError(f"Input directory does not exist: {input_dir}")

    images = collect_images(input_dir, tuple(args.extensions))
    if not images:
        raise FileNotFoundError(f"No supported images found in: {input_dir}")

    presentation = Presentation()
    presentation.slide_width = Inches(args.width)
    presentation.slide_height = Inches(args.height)
    presentation.core_properties.title = args.title or output.stem

    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    blank_layout = presentation.slide_layouts[6]

    for image_path in images:
        slide = presentation.slides.add_slide(blank_layout)
        add_background(slide, slide_width, slide_height, args.background)
        add_fitted_picture(slide, image_path, slide_width, slide_height, args.fit)

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)

    verification = Presentation(output)
    if len(verification.slides) != len(images):
        raise RuntimeError("Saved deck slide count does not match source image count")

    return {
        "output": str(output),
        "slides": len(images),
        "slide_size_inches": [args.width, args.height],
        "fit": args.fit,
        "images": [path.name for path in images],
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input_dir", type=Path, help="Directory containing ordered slide images")
    parser.add_argument("output", type=Path, help="Output .pptx path")
    parser.add_argument("--fit", choices=("cover", "contain", "stretch"), default="cover")
    parser.add_argument("--width", type=float, default=13.333333, help="Slide width in inches")
    parser.add_argument("--height", type=float, default=7.5, help="Slide height in inches")
    parser.add_argument("--background", default="FFFFFF", help="Six-digit RGB background color")
    parser.add_argument("--title", default="", help="PowerPoint document title")
    parser.add_argument("--extensions", nargs="+", default=list(DEFAULT_EXTENSIONS))
    args = parser.parse_args()

    if args.output.suffix.lower() != ".pptx":
        parser.error("output must use the .pptx extension")
    if args.width <= 0 or args.height <= 0:
        parser.error("slide width and height must be positive")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", args.background):
        parser.error("background must be a six-digit RGB value such as FFFFFF")

    args.background = args.background.upper()
    return args


def main() -> None:
    result = build_deck(parse_args())
    print(json.dumps(result, ensure_ascii=True, indent=2))


if __name__ == "__main__":
    main()
